import pandas as pd
import openpyxl
from pptx import Presentation
from tqdm import tqdm
import math
import os
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

class LRDM:
    def __init__(self, **args):
        self.input = "./input/"
        self.demographics_file = self.input + args['demographics']
        self.sample_file = self.input + "/templates/" + args['sample']
        self.analyst_file = self.input + args['analyst']
        self.GM_levels_file = self.input + args['gm_levels']
        self.leaders_file = self.input + args['leaders']

    def readAllFiles(self):
        self.demographics_pd = pd.read_excel(self.demographics_file, engine="openpyxl")

        self.analyst_leader_pd = pd.read_excel(self.analyst_file, engine="openpyxl", sheet_name="Leader")
        self.analyst_GM_pd = pd.read_excel(self.analyst_file, engine="openpyxl", sheet_name="GM")
        self.analyst_Site_pd = pd.read_excel(self.analyst_file, engine="openpyxl", sheet_name="Site Leader")

        self.GM_levels_pd = pd.read_excel(self.GM_levels_file, engine="openpyxl")

        self.leaders_pd = pd.read_excel(self.leaders_file, engine="openpyxl", sheet_name="Leader")
        self.GMs_pd = pd.read_excel(self.leaders_file, engine="openpyxl", sheet_name="GM")
        self.Site_leaders_pd = pd.read_excel(self.leaders_file, engine="openpyxl", sheet_name="Site Leader")

    def setLeader(self, id, GM=False, site_lead=False):
        self.leader_id = id
        self.GM = GM
        self.site_lead = site_lead

    def prepareContent(self):        
        _temp = "Supervisor Level {} ID"
        self.path = []

        _leader_entry = self.demographics_pd[self.demographics_pd.loc[:, "Worker ID"] == self.leader_id]
        _last_name = _leader_entry['Worker Last Name'].values[0]
        _full_name = _leader_entry['Worker Name'].values[0]
        self.first_content = [_full_name, _last_name]
        
        if self.GM:
            _temp_level = "GM Level {} ID"

            self.first_content = [self.GM[:-4], self.GM[:-4]]
            _analyst_entry = self.analyst_GM_pd[self.analyst_GM_pd["GM ID"] == self.leader_id]
            _gilead_entry = self.analyst_GM_pd[self.analyst_GM_pd["GM ID"] == 999999]
    
        elif self.site_lead:
            self.first_content = [self.site_lead.replace(" / ", " "), self.site_lead.replace(" / ", " ")]

        _output_path = './output/' + self.first_content[0] + '/'
        _filename = "2021-04 Global Employee Survey - {} - Leader Results Deck.pptx"
        _filename = _filename.format(self.first_content[1] + " Organization")

        self.path.append(_output_path + _filename)
        _source_path = './Deck Sample Folders/' + self.first_content[0] + '/'
        _filenames = [
            '2021-04 Global Employee Survey - {} - Score Summary.xlsx',
            '2021-04 Global Employee Survey - {} - Score Details.xlsx',
            '2021-04 Global Employee Survey - {} - Comment Themes.xlsx',
        ]
        [self.path.append(_source_path + _filename.format(self.first_content[1])) for _filename in _filenames]
        self.prs = Presentation(self.path[0])

    def makePresentation(self):
        shapes = self.prs.slides[3].shapes
        for i in range(3):
            shapes.add_embedded_xlsx(
                xlsx_file=self.path[i + 1], left=Inches(4 * i + 1), top=2202989, width=659686, height=1371600
            )
        print(self.path[0])

    def writeOutput(self):
        _output_path = './output/' + self.first_content[0] + '/'

        if os.path.exists(self.path[0]):
            os.remove(self.path[0])

        ## make a folder to involve the output file.
        os.makedirs(_output_path, exist_ok=True)

        self.prs.save(self.path[0])

if __name__ == "__main__":
    init_data = {
        'demographics': "2020 Demographics File Sample 2021-02-17.xlsx",
        'sample': "Leader Results Deck Master Template",
        'analyst': "Analyst Insights Database Sample 2021-02-25.xlsx",
        'gm_levels': "GM Levels 2021-02-17.xlsx",
        'leaders': "List of Leader GM SiteLeader 2021-02-23.xlsx",
    }

    lrdm = LRDM(**init_data)

    lrdm.readAllFiles()

    total_ids = len(lrdm.leaders_pd.index) + len(lrdm.GMs_pd.index) + len(lrdm.Site_leaders_pd.index)
    for index in tqdm(range(min(total_ids, 1)), desc="LRDM process"):
        if index < len(lrdm.leaders_pd.index):
            if lrdm.leaders_pd.iloc[index, :].values[0] == 999999:
                continue
            lrdm.setLeader(lrdm.leaders_pd.iloc[index, :].values[0])
        elif index < len(lrdm.leaders_pd.index) + len(lrdm.GMs_pd.index):
            row = lrdm.GMs_pd.iloc[index - len(lrdm.leaders_pd.index), :]
            lrdm.setLeader(row["GM ID"], row["GM Org"])
        else:
            row = lrdm.Site_leaders_pd.iloc[index - len(lrdm.leaders_pd.index) - len(lrdm.GMs_pd.index), :]
            lrdm.setLeader(row["Site Leader ID"], False, row["Site Name"])

        lrdm.prepareContent()

        lrdm.makePresentation()

        lrdm.writeOutput()
    print("done")